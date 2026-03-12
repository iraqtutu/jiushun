from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_ppt():
    prs = Presentation()

    # Define Corporate Colors (Green for Agriculture)
    PRIMARY_COLOR = RGBColor(0, 128, 0)  # Green
    SECONDARY_COLOR = RGBColor(105, 105, 105) # Dim Gray

    def add_title_slide(prs, title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        subtitle.text = subtitle_text
        
        # Style Title
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title.text_frame.paragraphs[0].font.bold = True
        
    def add_content_slide(prs, title_text, content_items):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Content
        content_place = slide.placeholders[1]
        tf = content_place.text_frame
        tf.clear() # Clear default bullet
        
        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.space_after = Pt(14)
            if item.startswith("-"):
                p.level = 1
            else:
                p.level = 0

    # 1. Cover Slide
    add_title_slide(prs, "玖顺农机售后服务工单系统", "用户使用手册 v1.0\n适用对象：玖顺员工、经销商、服务人员")

    # 2. System Intro
    add_content_slide(prs, "1. 系统简介", [
        "本系统专为玖顺农机售后服务体系打造，旨在实现：",
        "- 无纸化录入：告别纸质单据，手机端直接填写。",
        "- 数字化存储：数据云端保存，永不丢失。",
        "- 高效查询：随时随地查看历史维修记录。",
        "核心平台：微信小程序（无需下载安装APP）"
    ])

    # 3. Login & Register
    add_content_slide(prs, "2. 登录与账号申请", [
        "登录方式：",
        "- 打开“玖顺农机售后服务”小程序。",
        "- 点击“微信一键登录”即可。",
        "新用户申请：",
        "- 若首次登录，系统将引导进入申请页。",
        "- 填写姓名、手机号、选择角色（员工/经销商/服务人员）。",
        "- 提交后等待管理员审批通过即可使用。"
    ])

    # 4. Home Page
    add_content_slide(prs, "3. 首页功能一览", [
        "登录成功后进入首页，主要包含：",
        "1. 顶部信息卡：显示您的姓名及当前角色。",
        "2. 新建派工单：录入新的售后服务记录。",
        "3. 我的工单 / 工单查询：查看历史记录。",
        "4. 成员审批：(仅管理员) 审核新用户申请。"
    ])

    # 5. Create Work Order
    add_content_slide(prs, "4.1 新建派工单 (核心功能)", [
        "点击“+ 新建派工单”进入：",
        "- 基础信息：单号自动生成。",
        "- 产品信息：必填机器编号，必须上传铭牌照片。",
        "- 更换零件：每项零件可单独选择“带回”或“丢弃”。",
        "- 现场照片：最少1张，需与换件数量匹配。",
        "- 草稿功能：中途退出自动保存，随时恢复。"
    ])

    # 6. View Orders (Role Based)
    add_content_slide(prs, "4.2 查看与导出 (查询功能)", [
        "普通用户 (经销商/服务人员)：",
        "- 仅查看【自己提交】的工单，支持日期筛选。",
        "系统管理员 (Admin)：",
        "- 高级查询：按客户/电话/报单人/机型搜索。",
        "- 数据导出：点击右下角“导出汇总”按钮。",
        "  • 格式：Excel表格 (.xls)",
        "  • 内容：含所有工单详情及【照片高清链接】。"
    ])

    # 7. Admin Approval
    add_content_slide(prs, "4.3 成员审批 (管理员)", [
        "仅管理员可见此功能：",
        "- 待办列表：显示所有“待审核”的申请。",
        "- 通过：用户获得权限，可正常登录。",
        "- 拒绝：填写拒绝理由，用户需重新修改提交。",
        "提示：新用户注册后请及时联系管理员审核。"
    ])

    # 8. FAQ
    add_content_slide(prs, "5. 常见问题 (Q&A)", [
        "Q: 导出的Excel照片怎么看？",
        "A: 表格中包含照片链接，点击或复制到浏览器即可查看原图。",
        "Q: 填写中途退出数据会丢吗？",
        "A: 不会，系统有自动草稿功能，再次进入可恢复。",
        "Q: 图片上传失败？",
        "A: 可能是现场信号差，建议先拍照，到信号好的地方再上传。"
    ])

    # 9. End Slide
    add_title_slide(prs, "谢谢使用", "技术支持：系统开发维护团队\n遇到问题请截图反馈")

    # Save
    output_path = "Jiushun_User_Manual_v1.1.pptx"
    prs.save(output_path)
    print(f"Successfully created {output_path}")

if __name__ == "__main__":
    create_ppt()
